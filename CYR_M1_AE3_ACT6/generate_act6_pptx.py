"""
Generador de presentación PowerPoint para la Actividad N°6: Instalando Puntos de Red
Módulo 1 – Instalación de redes de área local cableadas e inalámbricas
Conectividad y Redes – 3° Medio TP
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colores ────────────────────────────────────────────────────────────────────
AZUL_OSCURO = RGBColor(0x1F, 0x38, 0x64)   # fondo
NARANJA     = RGBColor(0xFF, 0x66, 0x00)   # barra top / headers tabla
AZUL_MEDIO  = RGBColor(0x2E, 0x4B, 0x7A)   # fila alterna tabla
BLANCO      = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_CLARO  = RGBColor(0xD9, 0xD9, 0xD9)
VERDE       = RGBColor(0x70, 0xAD, 0x47)
ROJO        = RGBColor(0xFF, 0x00, 0x00)

# ── Dimensiones ────────────────────────────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
BAR_H   = Inches(0.25)


# ══════════════════════════════════════════════════════════════════════════════
# Helpers
# ══════════════════════════════════════════════════════════════════════════════

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def set_bg(slide, color=AZUL_OSCURO):
    """Rellena el fondo de la diapositiva con el color dado."""
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_bar(slide, color=NARANJA):
    """Barra naranja de 0.25" en la parte superior."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, SLIDE_W, BAR_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def add_textbox(slide, left, top, width, height, text,
                font_size=18, bold=False, color=BLANCO,
                align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_title(slide, title_text, subtitle=None):
    """Título principal centrado + subtítulo opcional."""
    add_textbox(slide,
                Inches(0.5), Inches(0.35),
                Inches(12.33), Inches(0.5),
                title_text,
                font_size=32, bold=True, color=BLANCO,
                align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide,
                    Inches(0.5), Inches(0.9),
                    Inches(12.33), Inches(0.35),
                    subtitle,
                    font_size=18, bold=False, color=NARANJA,
                    align=PP_ALIGN.LEFT)


def add_bullet_list(slide, items, top_start=Inches(1.35), left=Inches(0.5),
                    width=Inches(12.33), font_size=18, color=BLANCO,
                    line_height=Inches(0.42)):
    """Agrega una lista de bullets (texto plano)."""
    for i, item in enumerate(items):
        add_textbox(slide, left, top_start + i * line_height,
                    width, line_height, item,
                    font_size=font_size, color=color)


def add_table(slide, headers, rows,
              left=Inches(0.5), top=Inches(1.35),
              width=Inches(12.33), row_height=Inches(0.45)):
    """Tabla con header naranja y filas alternas azul oscuro / azul medio."""
    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 para header
    col_w  = width // n_cols

    tbl = slide.shapes.add_table(n_rows, n_cols,
                                  left, top,
                                  width, row_height * n_rows).table

    # ajustar alto de fila
    for r in range(n_rows):
        tbl.rows[r].height = row_height

    # Header
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.bold  = True
        run.font.size  = Pt(16)
        run.font.color.rgb = BLANCO
        cell.fill.solid()
        cell.fill.fore_color.rgb = NARANJA

    # Datos
    for r, row in enumerate(rows):
        bg = AZUL_MEDIO if r % 2 == 0 else AZUL_OSCURO
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size  = Pt(14)
            run.font.color.rgb = BLANCO
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    return tbl


# ══════════════════════════════════════════════════════════════════════════════
# Slides
# ══════════════════════════════════════════════════════════════════════════════

def slide_portada(prs):
    s = blank_slide(prs)
    set_bg(s)
    add_top_bar(s)

    # Banda naranja vertical decorativa
    band = s.shapes.add_shape(1, Inches(0), Inches(0.25), Inches(0.08), SLIDE_H - BAR_H)
    band.fill.solid(); band.fill.fore_color.rgb = NARANJA; band.line.fill.background()

    add_textbox(s, Inches(0.3), Inches(1.4), Inches(12.5), Inches(0.6),
                "Conectividad y Redes  |  Módulo 1",
                font_size=20, color=NARANJA, bold=True)

    add_textbox(s, Inches(0.3), Inches(2.2), Inches(12.5), Inches(1.1),
                "Actividad N°6:",
                font_size=36, bold=True, color=BLANCO)

    add_textbox(s, Inches(0.3), Inches(3.2), Inches(12.5), Inches(1.1),
                "Instalando Puntos de Red",
                font_size=44, bold=True, color=NARANJA)

    add_textbox(s, Inches(0.3), Inches(4.5), Inches(12.5), Inches(0.5),
                "Instalación de redes de área local cableadas e inalámbricas",
                font_size=20, color=BLANCO)

    add_textbox(s, Inches(0.3), Inches(5.2), Inches(12.5), Inches(0.5),
                "3° Medio  ·  Técnico Profesional  ·  Especialidad Conectividad y Redes",
                font_size=16, color=GRIS_CLARO, italic=True)


def slide_objetivos(prs):
    s = blank_slide(prs)
    set_bg(s)
    add_top_bar(s)
    add_title(s, "Objetivos de la Actividad",
              "Aprendizaje Esperado 1.3")

    criterios = [
        "✅  1.3.3 – Clasificar materiales, herramientas e implementos de seguridad para instalar cableado estructurado.",
        "✅  1.3.4 – Instalar ductos y cableado según planos; etiquetar cables, conectores y tomas de red.",
        "✅  1.3.5 – Instalar racks, patch panels y puntos terminales de red respetando estándares técnicos.",
        "✅  1.3.6 – Realizar pruebas con LAN Tester para verificar la calidad de la instalación.",
        "✅  1.3.7 – Comunicar e informar los resultados generando documentación pertinente.",
    ]
    add_bullet_list(s, criterios, top_start=Inches(1.4),
                    line_height=Inches(0.9), font_size=17)


def slide_que_es_punto_red(prs):
    s = blank_slide(prs)
    set_bg(s)
    add_top_bar(s)
    add_title(s, "¿Qué es un Punto de Red?",
              "La conexión final entre el cableado y el usuario")

    add_textbox(s, Inches(0.5), Inches(1.35), Inches(7.5), Inches(1.5),
                "Un punto de red es la toma o salida donde el usuario conecta "
                "su computador a la red de cableado estructurado del edificio.\n"
                "Se compone de: Jack Keystone + Faceplate + Canaleta.",
                font_size=18, color=BLANCO)

    add_textbox(s, Inches(0.5), Inches(3.0), Inches(4.0), Inches(0.4),
                "💡 Ejemplo cotidiano:", font_size=18, bold=True, color=NARANJA)

    add_textbox(s, Inches(0.5), Inches(3.5), Inches(12.0), Inches(1.5),
                "Imagina que la red del colegio es como el tendido eléctrico de una casa:\n"
                "el cable UTP sería el cable eléctrico dentro de la pared,\n"
                "y el punto de red sería el enchufe en la pared que ves en cada sala.",
                font_size=18, color=BLANCO)

    # Recuadro visual
    box = s.shapes.add_shape(1, Inches(8.5), Inches(1.4), Inches(4.2), Inches(4.5))
    box.fill.solid(); box.fill.fore_color.rgb = AZUL_MEDIO
    box.line.color.rgb = NARANJA

    add_textbox(s, Inches(8.6), Inches(1.5), Inches(4.0), Inches(4.0),
                "🖥️  PC del alumno\n         |\n🔌  Patch Cord\n         |\n🔲  Jack Keystone (pared)\n         |\n📦  Cable UTP Cat6 (canaleta)\n         |\n🗂️  Patch Panel\n         |\n🔀  Switch → 🌐 Internet",
                font_size=14, color=BLANCO)


def slide_materiales(prs):
    s = blank_slide(prs)
    set_bg(s)
    add_top_bar(s)
    add_title(s, "Materiales y Herramientas Necesarias",
              "¿Qué necesitamos para instalar un punto de red?")

    headers = ["Material / Herramienta", "¿Para qué sirve?"]
    rows = [
        ("Cable UTP Cat6",            "Transmitir datos hasta 100 m"),
        ("Jack Keystone RJ-45",        "Toma de red empotrable en la pared"),
        ("Faceplate",                  "Cubre el jack; da terminación prolija"),
        ("Canaleta plástica",          "Protege y ordena los cables en la pared"),
        ("Patch Panel 24 puertos",     "Organiza todas las llegadas de cable al rack"),
        ("Herramienta Punch-Down",     "Presiona los cables en el jack o patch panel"),
        ("Tijeras / Pelacable",        "Quitar la cubierta exterior del cable"),
        ("LAN Tester",                 "Probar que el cable está bien terminado"),
        ("Etiquetas adhesivas",        "Identificar cada cable (estándar TIA)"),
    ]
    add_table(s, headers, rows,
              top=Inches(1.35), row_height=Inches(0.52))


def slide_cable_utp(prs):
    s = blank_slide(prs)
    set_bg(s)
    add_top_bar(s)
    add_title(s, "Cable UTP Cat6 – Características",
              "El 'conductor' de los datos dentro del edificio")

    items = [
        "📏  Distancia máxima: 100 metros por segmento (estándar TIA-568)",
        "🔢  8 conductores de cobre organizados en 4 pares trenzados",
        "⚡  Velocidad: hasta 10 Gbps (Cat6) – Cat5e solo llega a 1 Gbps",
        "🎨  Colores de pares: Azul / Naranja / Verde / Marrón (+ variante blanca)",
        "🔒  Sin blindaje en UTP (Unshielded Twisted Pair) – más económico",
        "🏭  Estándar: ANSI/TIA-568-C.2",
    ]
    add_bullet_list(s, items, top_start=Inches(1.4),
                    line_height=Inches(0.75), font_size=18)

    add_textbox(s, Inches(0.5), Inches(5.7), Inches(12.0), Inches(0.5),
                "💡 Cat6 tiene un separador plástico central (cruceta) que reduce la interferencia entre pares.",
                font_size=16, color=NARANJA, italic=True)


def slide_jack_keystone(prs):
    s = blank_slide(prs)
    set_bg(s)
    add_top_bar(s)
    add_title(s, "Jack Keystone – ¿Qué es y cómo funciona?",
              "La pieza clave del punto de red")

    add_textbox(s, Inches(0.5), Inches(1.35), Inches(7.8), Inches(2.0),
                "El Jack Keystone es un módulo RJ-45 que se encaja (snap-in) "
                "en un faceplate o en el patch panel.\n\n"
                "▶  Parte trasera: 8 ranuras donde se insertan los cables con la herramienta punch-down\n"
                "▶  Parte frontal: conector RJ-45 hembra donde el usuario conecta su patch cord",
                font_size=17, color=BLANCO)

    # Cuadro tipo recuerda
    box = s.shapes.add_shape(1, Inches(0.5), Inches(3.5), Inches(12.0), Inches(2.5))
    box.fill.solid(); box.fill.fore_color.rgb = AZUL_MEDIO
    box.line.color.rgb = NARANJA

    add_textbox(s, Inches(0.7), Inches(3.6), Inches(11.5), Inches(2.2),
                "📌 Recuerda:\n"
                "Los jacks son de colores o tienen etiquetas A/B para indicar la norma (568A o 568B).\n"
                "Siempre usa la misma norma en AMBOS extremos del cable → si empiezas con 568B, termina con 568B.\n"
                "Mezclar normas en los dos extremos crea un cable CRUZADO (útil para otros casos).",
                font_size=16, color=BLANCO)


def slide_tia568b(prs):
    s = blank_slide(prs)
    set_bg(s)
    add_top_bar(s)
    add_title(s, "Norma TIA 568B – Tabla de Colores",
              "La norma más usada en Chile para instalaciones comerciales")

    headers = ["Pin", "Color", "Par", "Función"]
    rows = [
        ("1", "Blanco-Naranja",  "2", "TX+ (transmisión)"),
        ("2", "Naranja",         "2", "TX- (transmisión)"),
        ("3", "Blanco-Verde",    "3", "RX+ (recepción)"),
        ("4", "Azul",            "1", "PoE / telefonía"),
        ("5", "Blanco-Azul",     "1", "PoE / telefonía"),
        ("6", "Verde",           "3", "RX- (recepción)"),
        ("7", "Blanco-Marrón",   "4", "PoE / Gigabit"),
        ("8", "Marrón",          "4", "PoE / Gigabit"),
    ]
    add_table(s, headers, rows, top=Inches(1.3), row_height=Inches(0.5))

    add_textbox(s, Inches(0.5), Inches(6.1), Inches(12.0), Inches(0.5),
                "💡 TIA 568A invierte pares 2 y 3: Pin 1=Blanco-Verde, 2=Verde, 3=Blanco-Naranja, 6=Naranja",
                font_size=14, color=NARANJA, italic=True)


def slide_punch_down(prs):
    s = blank_slide(prs)
    set_bg(s)
    add_top_bar(s)
    add_title(s, "Herramienta Punch-Down – ¿Cómo usarla?",
              "La 'pistola' del instalador de red")

    add_textbox(s, Inches(0.5), Inches(1.35), Inches(7.5), Inches(1.5),
                "La herramienta punch-down (impactadora) inserta y corta el cable conductor "
                "dentro de la ranura del jack o patch panel en un solo golpe.",
                font_size=18, color=BLANCO)

    pasos = [
        "1️⃣  Quitar aprox. 3 cm de cubierta exterior del cable sin dañar los pares.",
        "2️⃣  Separar los 4 pares y ordenarlos según la norma del jack (A o B).",
        "3️⃣  Insertar cada conductor en su ranura de color correspondiente.",
        "4️⃣  Colocar la herramienta punch-down con el filo hacia AFUERA del jack.",
        "5️⃣  Presionar firmemente hasta escuchar un 'click' – el cable queda sujeto y cortado.",
        "⚠️  No destrenzes más de 1.25 cm de los pares – ¡afecta la calidad de la señal!",
    ]
    add_bullet_list(s, pasos, top_start=Inches(3.1),
                    line_height=Inches(0.62), font_size=17)


def slide_paso_a_paso(prs):
    s = blank_slide(prs)
    set_bg(s)
    add_top_bar(s)
    add_title(s, "Paso a Paso: Instalación del Jack Keystone",
              "6 pasos para conectar un punto de red correctamente")

    pasos = [
        "1.  Mide y corta el cable UTP Cat6 dejando 30 cm de margen en la caja.",
        "2.  Pela 3 cm de la cubierta exterior con el pelacable o tijera.",
        "3.  Ordena los 8 conductores según la norma TIA 568B (colores del jack).",
        "4.  Usa el punch-down para presionar cada cable en su ranura.",
        "5.  Ajusta el jack al faceplate y fija el faceplate a la caja.",
        "6.  Prueba con el LAN Tester – los 8 LEDs deben parpadear en orden.",
    ]
    add_bullet_list(s, pasos, top_start=Inches(1.45),
                    line_height=Inches(0.8), font_size=18)


def slide_faceplate_canaleta(prs):
    s = blank_slide(prs)
    set_bg(s)
    add_top_bar(s)
    add_title(s, "Faceplate y Canaleta – Para qué sirven",
              "La 'terminación' profesional del punto de red")

    col_l = Inches(0.5)
    col_r = Inches(6.8)
    w_col = Inches(5.8)

    # Columna izquierda - Faceplate
    add_textbox(s, col_l, Inches(1.35), w_col, Inches(0.4),
                "🔲  Faceplate (placa de pared)", font_size=20, bold=True, color=NARANJA)
    items_l = [
        "• Marco plástico que cubre la caja empotrada.",
        "• Acepta 1, 2 o 4 módulos keystone.",
        "• Da terminación estética y profesional.",
        "• Protege el jack del polvo y golpes.",
    ]
    add_bullet_list(s, items_l, top_start=Inches(1.85), left=col_l,
                    width=w_col, line_height=Inches(0.58), font_size=16)

    # Columna derecha - Canaleta
    add_textbox(s, col_r, Inches(1.35), w_col, Inches(0.4),
                "📦  Canaleta de red", font_size=20, bold=True, color=NARANJA)
    items_r = [
        "• Canal plástico que recorre paredes/cielo.",
        "• Protege y ordena los cables UTP.",
        "• Ancho varía según N° de cables (20x12 a 100x60 mm).",
        "• Cumple norma de cableado estructurado TIA-568.",
    ]
    add_bullet_list(s, items_r, top_start=Inches(1.85), left=col_r,
                    width=w_col, line_height=Inches(0.58), font_size=16)

    add_textbox(s, Inches(0.5), Inches(5.5), Inches(12.0), Inches(0.65),
                "💡 Para 24 puntos de red, una canaleta 40×25 mm suele ser suficiente. Para más puntos, usar 60×40 mm.",
                font_size=16, color=NARANJA, italic=True)


def slide_patch_panel(prs):
    s = blank_slide(prs)
    set_bg(s)
    add_top_bar(s)
    add_title(s, "Conexión al Patch Panel – El Recorrido Completo",
              "Desde el computador del alumno hasta el switch")

    add_textbox(s, Inches(0.5), Inches(1.35), Inches(12.0), Inches(0.5),
                "El recorrido de la señal en una red con cableado estructurado:", font_size=18, color=NARANJA, bold=True)

    steps = [
        ("🖥️  PC / Laptop",         "Genera el tráfico de datos"),
        ("🔌  Patch Cord (usuario)", "Cable corto RJ45 ↔ Jack de pared"),
        ("🔲  Jack Keystone",         "Toma de red empotrada en la pared"),
        ("📦  Cable UTP Cat6",        "Recorre la canaleta hasta el rack (máx. 90 m)"),
        ("🗂️  Patch Panel",           "Organiza todas las llegadas de cable en el rack"),
        ("🔌  Patch Cord (equipo)",  "Cable corto Patch Panel ↔ Switch"),
        ("🔀  Switch",                "Conecta todos los equipos de la red"),
        ("🌐  Router / Internet",    "Salida hacia la red exterior"),
    ]

    for i, (comp, desc) in enumerate(steps):
        y = Inches(1.95) + i * Inches(0.58)
        bg = AZUL_MEDIO if i % 2 == 0 else AZUL_OSCURO
        box = s.shapes.add_shape(1, Inches(0.4), y, Inches(12.5), Inches(0.54))
        box.fill.solid(); box.fill.fore_color.rgb = bg; box.line.fill.background()
        add_textbox(s, Inches(0.5), y + Inches(0.05),
                    Inches(4.0), Inches(0.45), comp, font_size=15, bold=True, color=NARANJA)
        add_textbox(s, Inches(4.7), y + Inches(0.05),
                    Inches(8.0), Inches(0.45), desc, font_size=15, color=BLANCO)


def slide_errores_comunes(prs):
    s = blank_slide(prs)
    set_bg(s)
    add_top_bar(s)
    add_title(s, "Errores Comunes – ¡Qué NO Hacer!",
              "Aprende de los errores más frecuentes en instalación")

    errores = [
        "❌  Mezclar norma 568A en un extremo y 568B en el otro (sin querer).",
        "❌  Pelar demasiado el cable → conductores expuestos = interferencia.",
        "❌  Destrensar más de 1.25 cm los pares → pérdida de señal a alta velocidad.",
        "❌  Doblar el cable en ángulo de 90° → aplasta los conductores y rompe el par.",
        "❌  No etiquetar el cable → confusión futura y tiempo perdido en diagnóstico.",
        "❌  Superar los 100 m de cable → la señal se degrada y el tester fallará.",
        "❌  No testear antes de cerrar la pared → ¡descubierto demasiado tarde!",
    ]
    add_bullet_list(s, errores, top_start=Inches(1.45),
                    line_height=Inches(0.71), font_size=17)


def slide_tester_red(prs):
    s = blank_slide(prs)
    set_bg(s)
    add_top_bar(s)
    add_title(s, "Cómo Probar el Punto de Red con LAN Tester",
              "La prueba de calidad antes de dar el trabajo por terminado")

    add_textbox(s, Inches(0.5), Inches(1.35), Inches(12.0), Inches(0.5),
                "El LAN Tester mide la CONTINUIDAD de cada uno de los 8 pines del cable.", font_size=18, color=BLANCO)

    headers = ["LED N°", "Pin / Color (568B)", "¿Qué indica?"]
    rows = [
        ("1", "Blanco-Naranja",  "Parpadea → OK  /  No parpadea → corte"),
        ("2", "Naranja",         "Parpadea → OK  /  No parpadea → corte"),
        ("3", "Blanco-Verde",    "Parpadea → OK  /  No parpadea → corte"),
        ("4", "Azul",            "Parpadea → OK  /  No parpadea → corte"),
        ("5", "Blanco-Azul",     "Parpadea → OK  /  No parpadea → corte"),
        ("6", "Verde",           "Parpadea → OK  /  No parpadea → corte"),
        ("7", "Blanco-Marrón",   "Parpadea → OK  /  No parpadea → corte"),
        ("8", "Marrón",          "Parpadea → OK  /  No parpadea → corte"),
    ]
    add_table(s, headers, rows, top=Inches(2.0), row_height=Inches(0.48))

    add_textbox(s, Inches(0.5), Inches(6.5), Inches(12.0), Inches(0.5),
                "✅ Prueba exitosa: los 8 LEDs parpadean en orden (1-2-3-4-5-6-7-8). Si parpadean fuera de orden → cable cruzado.",
                font_size=15, color=NARANJA, italic=True)


def slide_instrucciones_actividad(prs):
    s = blank_slide(prs)
    set_bg(s)
    add_top_bar(s)
    add_title(s, "Instrucciones de la Actividad",
              "Trabajo en clases – Realizando mis puntos de red")

    instrucciones = [
        "1.  Formen equipos de 3 integrantes.",
        "2.  Armen un cable 'cruzado' y un cable 'derecho' con conectores RJ-45\n     siguiendo normas TIA 568A y TIA 568B.",
        "3.  Muestren el cable al docente y compartan resultados con el curso.",
        "4.  Verifiquen el mapeado y continuidad del cable con el probador de cables.",
        "5.  Etiqueten ambos extremos según estándares internacionales.",
        "6.  Conecten los patch cords entre switch y patch panel del bastidor.",
        "7.  Dejen el lugar limpio y ordenado — ¡el orden es parte de la nota!",
    ]
    add_bullet_list(s, instrucciones, top_start=Inches(1.4),
                    line_height=Inches(0.75), font_size=17)


def slide_estudio_de_caso(prs):
    s = blank_slide(prs)
    set_bg(s)
    add_top_bar(s)
    add_title(s, "Estudio de Caso",
              "Situación real: laboratorio de computación")

    add_textbox(s, Inches(0.5), Inches(1.35), Inches(12.0), Inches(1.6),
                "Un establecimiento educacional ha adjudicado a tu empresa la instalación de "
                "una red de cableado estructurado horizontal para abastecer un laboratorio "
                "de computación.\n"
                "La red se expande desde el armario/gabinete central hasta las áreas de trabajo "
                "de cada usuario.",
                font_size=18, color=BLANCO)

    add_textbox(s, Inches(0.5), Inches(3.1), Inches(12.0), Inches(0.4),
                "📋  Tareas asignadas:", font_size=18, bold=True, color=NARANJA)

    tareas = [
        "→  Montar switch y patch panel en el armario respetando 2U de separación.",
        "→  Construir 2 patch cords según normativa TIA/EIA-568.",
        "→  Instalar ductos y canaletas siguiendo el plano de red.",
        "→  Conectar módulos RJ-45 (jacks) y user cords en cada punto de trabajo.",
        "→  Testar y etiquetar cada punto de red al finalizar.",
    ]
    add_bullet_list(s, tareas, top_start=Inches(3.6),
                    line_height=Inches(0.66), font_size=17)


def slide_rack_armario(prs):
    s = blank_slide(prs)
    set_bg(s)
    add_top_bar(s)
    add_title(s, "Rack / Armario – El corazón de la red",
              "Donde se concentra todo el equipamiento")

    col_l = Inches(0.5)
    col_r = Inches(6.8)
    w_col = Inches(5.8)

    add_textbox(s, col_l, Inches(1.35), w_col, Inches(0.4),
                "📦  ¿Qué es?", font_size=20, bold=True, color=NARANJA)
    add_textbox(s, col_l, Inches(1.8), w_col, Inches(2.0),
                "Gabinete metálico de 19\" (estándar) que contiene:\n"
                "• Switch(es)\n• Patch panel\n• Bandejas de cables\n• UPS (respaldo eléctrico)",
                font_size=17, color=BLANCO)

    add_textbox(s, col_r, Inches(1.35), w_col, Inches(0.4),
                "📏  Medidas importantes", font_size=20, bold=True, color=NARANJA)
    add_textbox(s, col_r, Inches(1.8), w_col, Inches(2.0),
                "• 1 Unidad de Rack (1U) = 44,45 mm de alto\n"
                "• Patch panel típico ocupa 1U\n"
                "• Switch 24 puertos ocupa 1U\n"
                "• Separación recomendada: 2U entre dispositivos",
                font_size=17, color=BLANCO)

    add_textbox(s, Inches(0.5), Inches(4.1), Inches(12.0), Inches(0.4),
                "🔢  ¿Cuántos U necesito?", font_size=18, bold=True, color=NARANJA)
    add_textbox(s, Inches(0.5), Inches(4.55), Inches(12.0), Inches(0.9),
                "Ejemplo para laboratorio con 24 PCs: 1U patch panel + 2U separación + 1U switch + 2U separación "
                "+ 1U switch extra + 2U bandejas = aprox. 12U → usar rack de 16U o 22U.",
                font_size=17, color=BLANCO)


def slide_etiquetado(prs):
    s = blank_slide(prs)
    set_bg(s)
    add_top_bar(s)
    add_title(s, "Etiquetado de Cables – Norma TIA-606",
              "¿Por qué etiquetar? Porque sin orden no hay red que funcione bien")

    add_textbox(s, Inches(0.5), Inches(1.35), Inches(12.0), Inches(1.0),
                "El estándar TIA-606 establece cómo identificar cada cable, puerto y dispositivo "
                "en una instalación de cableado estructurado. Un buen etiquetado ahorra horas "
                "de trabajo en diagnóstico y mantenimiento.",
                font_size=18, color=BLANCO)

    headers = ["Elemento", "Formato etiqueta", "Ejemplo"]
    rows = [
        ("Cable de pared",    "SALA-PUNTO",          "LAB1-P03"),
        ("Puerto patch panel", "PP-SALA-N°",          "PP-LAB1-03"),
        ("Puerto switch",      "SW-N°",               "SW-03"),
        ("Jack de pared",      "SALA-PUNTO",          "LAB1-P03"),
        ("Patch cord",         "ORIGEN→DESTINO",      "SW03→PP03"),
    ]
    add_table(s, headers, rows, top=Inches(2.5), row_height=Inches(0.52))

    add_textbox(s, Inches(0.5), Inches(5.7), Inches(12.0), Inches(0.5),
                "⚠️  Etiquetar AMBOS extremos del cable. Usar cinta de etiquetado resistente o portaetiquetas.",
                font_size=15, color=NARANJA, italic=True)


def slide_entregables_evaluacion(prs):
    s = blank_slide(prs)
    set_bg(s)
    add_top_bar(s)
    add_title(s, "Entregables y Criterios de Evaluación",
              "Rúbrica de evaluación – Presentación de solución del caso")

    add_textbox(s, Inches(0.5), Inches(1.35), Inches(12.0), Inches(0.4),
                "📋  Tu presentación debe incluir:", font_size=18, bold=True, color=NARANJA)

    entregables = [
        "a.  Argumentación técnica paso a paso de los procedimientos realizados.",
        "b.  Descripción de infraestructura, herramientas y materiales utilizados.",
        "c.  Evidencia gráfica (fotos o diagramas) de los procedimientos.",
        "d.  Conclusión técnica grupal sobre el caso planteado.",
    ]
    add_bullet_list(s, entregables, top_start=Inches(1.85),
                    line_height=Inches(0.55), font_size=16)

    headers = ["Indicador",  "Excelente (4)", "Bueno (3)", "Regular (2)", "Debe mejorar (1)"]
    rows = [
        ("Argumentación técnica procedimientos", "100% de pasos argumentados", "75% argumentados", "50% argumentados", "≤25% argumentados"),
        ("Herramientas y materiales",            "Todos correctamente descritos", "La mayoría correctos", "Solo algunos", "Descripción incorrecta"),
        ("Evidencia gráfica",                    "Fotos claras de todos los pasos", "Fotos parciales", "Pocas evidencias", "Sin evidencias"),
        ("Conclusión técnica",                   "Completa y fundamentada", "Parcialmente fundamentada", "Superficial", "Sin conclusión"),
    ]
    add_table(s, headers, rows, top=Inches(3.5), row_height=Inches(0.55))


def slide_ticket_salida(prs):
    s = blank_slide(prs)
    set_bg(s)
    add_top_bar(s)
    add_title(s, "Ticket de Salida – ¡Reflexionemos!",
              "Antes de salir, responde estas preguntas")

    preguntas = [
        "1.  Para una red de 24 puntos de red, ¿qué canaleta recomendarías? ¿Por qué?",
        "2.  ¿Qué diferencia física tiene un cable UTP Cat6 vs Cat5e?",
        "3.  ¿Qué procedimiento te resultó más complejo en esta actividad?",
        "4.  ¿En qué situaciones de la vida real aplicarías estos conocimientos?",
        "5.  ¿Crees que podrías crear un emprendimiento con estas habilidades?",
    ]
    add_bullet_list(s, preguntas, top_start=Inches(1.45),
                    line_height=Inches(0.92), font_size=17)

    add_textbox(s, Inches(0.5), Inches(6.0), Inches(12.0), Inches(0.5),
                "💬  Puedes responder en Kahoot, Mentimeter, Padlet o Google Forms.",
                font_size=15, color=NARANJA, italic=True)


def slide_resumen(prs):
    s = blank_slide(prs)
    set_bg(s)
    add_top_bar(s)
    add_title(s, "Resumen – Puntos Clave Aprendidos",
              "¿Qué nos llevamos de esta actividad?")

    puntos = [
        "🔵  Un punto de red = Jack + Faceplate + Canaleta + Cable UTP Cat6",
        "🔵  Norma TIA 568B: el estándar más usado en instalaciones comerciales",
        "🔵  Punch-down: presionar y cortar en un golpe. Destrensar máximo 1,25 cm",
        "🔵  Distancia máxima: 100 m por segmento (90 m fijo + 10 m patch cords)",
        "🔵  Patch Panel: central de organización de cables en el rack",
        "🔵  LAN Tester: 8 LEDs en orden → cable OK, fuera de orden → cruzado",
        "🔵  Etiquetar SIEMPRE con la norma TIA-606 en ambos extremos",
        "🔵  Rack: gabinete 19\", altura en U (1U = 44,45 mm), separar 2U entre equipos",
    ]
    add_bullet_list(s, puntos, top_start=Inches(1.4),
                    line_height=Inches(0.68), font_size=17)


# ══════════════════════════════════════════════════════════════════════════════
# Main
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    slide_portada(prs)               # 1
    slide_objetivos(prs)             # 2
    slide_que_es_punto_red(prs)      # 3
    slide_materiales(prs)            # 4
    slide_cable_utp(prs)             # 5
    slide_jack_keystone(prs)         # 6
    slide_tia568b(prs)               # 7
    slide_punch_down(prs)            # 8
    slide_paso_a_paso(prs)           # 9
    slide_faceplate_canaleta(prs)    # 10
    slide_patch_panel(prs)           # 11
    slide_rack_armario(prs)          # 12
    slide_errores_comunes(prs)       # 13
    slide_tester_red(prs)            # 14
    slide_etiquetado(prs)            # 15
    slide_estudio_de_caso(prs)       # 16
    slide_instrucciones_actividad(prs)  # 17
    slide_entregables_evaluacion(prs)   # 18
    slide_ticket_salida(prs)         # 19
    slide_resumen(prs)               # 20

    output = "actividad6_puntos_de_red.pptx"
    prs.save(output)
    print(f"✅  Presentación generada: {output}  ({len(prs.slides)} diapositivas)")


if __name__ == "__main__":
    main()
