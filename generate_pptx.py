"""
generate_pptx.py
Genera resumen_modulos.pptx con 31 diapositivas para
Instalación de Redes – 3° Medio TP Conectividad y Redes.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colores ──────────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1F, 0x38, 0x64)   # #1F3864
ORANGE      = RGBColor(0xFF, 0x66, 0x00)   # #FF6600
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
MID_BLUE    = RGBColor(0x2E, 0x4B, 0x7A)   # fila alterna tabla
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)

# ── Presentación ─────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # layout en blanco


# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)


def set_background(slide, color=DARK_BLUE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def orange_bar(slide, height=Inches(0.35)):
    """Barra decorativa naranja en la parte superior."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0),
        prs.slide_width, height,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()


def add_title(slide, text, top=Inches(0.4), font_size=Pt(34)):
    txb = slide.shapes.add_textbox(
        Inches(0.4), top, Inches(12.5), Inches(0.9)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"


def add_body(slide, lines, top=Inches(1.5), font_size=Pt(18)):
    """Agrega un bloque de texto con viñetas."""
    txb = slide.shapes.add_textbox(
        Inches(0.5), top, Inches(12.3), Inches(5.5)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.space_after = Pt(4)


def add_table(slide, headers, rows, top=Inches(1.5), col_widths=None):
    """Crea una tabla con encabezado naranja y filas alternas."""
    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 para encabezado
    left   = Inches(0.4)
    width  = Inches(12.5)
    height = Inches(0.4 * n_rows + 0.1)

    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    # Anchos de columna
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)

    def set_cell(cell, text, bg, bold=False, font_size=Pt(15)):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.bold = bold
        p.font.size = font_size
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.LEFT

    # Encabezado
    for j, h in enumerate(headers):
        set_cell(tbl.cell(0, j), h, ORANGE, bold=True)

    # Filas
    for i, row in enumerate(rows):
        bg = MID_BLUE if i % 2 == 0 else DARK_BLUE
        for j, val in enumerate(row):
            set_cell(tbl.cell(i + 1, j), val, bg)


# =============================================================================
# SLIDE 1 – Portada
# =============================================================================
sl = add_slide()
set_background(sl)

# Rectángulo grande central con naranja
rect = sl.shapes.add_shape(1, Inches(1), Inches(1.5), Inches(11.33), Inches(4.5))
rect.fill.solid()
rect.fill.fore_color.rgb = MID_BLUE
rect.line.color.rgb = ORANGE
rect.line.width = Pt(3)

txb = sl.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11), Inches(1.5))
tf = txb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Resumen de Módulos – Instalación de Redes"
p.font.bold = True
p.font.size = Pt(40)
p.font.color.rgb = WHITE
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

txb2 = sl.shapes.add_textbox(Inches(1.2), Inches(3.5), Inches(11), Inches(1.0))
tf2 = txb2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "3° Medio TP Conectividad y Redes"
p2.font.size = Pt(26)
p2.font.color.rgb = ORANGE
p2.font.name = "Calibri"
p2.alignment = PP_ALIGN.CENTER

txb3 = sl.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(11), Inches(0.7))
tf3 = txb3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "Instalación de Redes – Asignatura Técnico-Profesional"
p3.font.size = Pt(18)
p3.font.color.rgb = LIGHT_GRAY
p3.font.name = "Calibri"
p3.alignment = PP_ALIGN.CENTER

# =============================================================================
# SLIDE 2 – Índice
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "📋 Índice de Contenidos")
add_body(sl, [
    "🔵  TEMA 1 – Modelo OSI y Modelo TCP/IP  (Slides 3–10)",
    "",
    "🔌  TEMA 2 – Componentes de Red: Router, Switch, Patch Panel,",
    "           Jack/Wall Mount, Rack, Servidores  (Slides 11–18)",
    "",
    "🔧  TEMA 3 – Cableado Estructurado: backbone, salas, acometida,",
    "           marcado de cables, subsistemas  (Slides 19–25)",
    "",
    "📏  TEMA 4 – Normas TIA/EIA 568A y 568B: colores, cable directo",
    "           vs cruzado, buenas prácticas  (Slides 26–31)",
])

# =============================================================================
# SLIDE 3 – ¿Qué es el Modelo OSI?
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "🔵 ¿Qué es el Modelo OSI?")
add_body(sl, [
    "📌 Definición: Estándar internacional de 7 capas para comunicaciones entre computadores.",
    "",
    "🗓️  Creado en 1984 por la ISO (International Organization for Standardization).",
    "",
    "🎯 Objetivo: Que equipos de DISTINTOS fabricantes puedan comunicarse entre sí.",
    "",
    "📐 OSI = Open Systems Interconnection (Interconexión de Sistemas Abiertos).",
    "",
    "✅ Modelo TEÓRICO de referencia — base para entender cómo funciona cualquier red.",
    "",
    '💡 Ejemplo: "Gracias al modelo OSI, un PC Dell puede hablar con un Mac de Apple."',
])

# =============================================================================
# SLIDE 4 – Las 7 Capas del Modelo OSI
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "📊 Las 7 Capas del Modelo OSI")
add_table(sl,
    ["N°", "Nombre", "Función", "Ejemplo"],
    [
        ["7", "Aplicación",      "Interfaz usuario ↔ red",        "HTTP, FTP, DNS, SMTP"],
        ["6", "Presentación",    "Cifrado y compresión",           "SSL/TLS, JPEG, MP4"],
        ["5", "Sesión",          "Abrir / cerrar sesiones",        "NetBIOS, RPC"],
        ["4", "Transporte",      "TCP/UDP, control de flujo",      "TCP puerto 80 (HTTP)"],
        ["3", "Red",             "Enrutamiento IP",                "Router, IPv4/IPv6"],
        ["2", "Enlace de Datos", "Tramas, dirección MAC",          "Switch, NIC"],
        ["1", "Física",          "Bits eléctricos / ópticos",      "Cable UTP, fibra óptica"],
    ],
    top=Inches(1.4),
    col_widths=[0.5, 1.8, 4.2, 5.5],
)

# =============================================================================
# SLIDE 5 – Capa Física (Capa 1)
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "⚡ Capa Física – Capa 1")
add_body(sl, [
    "❓ ¿Qué hace?: Convierte datos en señales eléctricas, ópticas o inalámbricas.",
    "",
    "📡 Medios de transmisión:",
    "    • Cable UTP Cat6 → señal eléctrica",
    "    • Fibra Óptica   → señal de luz (pulsos láser)",
    "    • WiFi           → ondas de radio (2.4 GHz / 5 GHz)",
    "",
    "🔧 Dispositivos de Capa 1: Hub, repetidor, cable de red.",
    "",
    '💡 Ejemplo real: "El cable de red que conectas al PC trabaja en Capa 1 — solo transporta bits, no sabe qué contienen."',
])

# =============================================================================
# SLIDE 6 – Capa de Red (Capa 3)
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "🌐 Capa de Red – Capa 3")
add_body(sl, [
    "❓ ¿Qué hace?: Enruta paquetes entre redes DISTINTAS.",
    "",
    "📦 Protocolo principal: IP (Internet Protocol)",
    "    • IPv4: direcciones de 32 bits  →  ej. 192.168.1.1",
    "    • IPv6: direcciones de 128 bits →  ej. 2001:db8::1",
    "",
    "🔧 Dispositivo de Capa 3: ROUTER",
    "",
    "🔄 El router examina la IP de destino y decide la mejor ruta.",
    "",
    '💡 Ejemplo: "Tu router de casa decide si el paquete va a Internet (WAN) o se queda en tu red local (LAN)."',
])

# =============================================================================
# SLIDE 7 – Capa de Transporte (Capa 4) – TCP vs UDP
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "🚦 Capa de Transporte – Capa 4: TCP vs UDP")
add_table(sl,
    ["Característica", "TCP", "UDP"],
    [
        ["Confiabilidad",  "✅ Garantiza entrega",      "❌ Sin confirmación"],
        ["Velocidad",      "Más lento (confirmaciones)", "Más rápido"],
        ["Orden",          "Ordena los paquetes",        "Sin orden garantizado"],
        ["Uso típico",     "Web, email, FTP",            "Video streaming, juegos, DNS"],
        ["Ejemplo",        "Descarga de archivo",        "Videollamada en tiempo real"],
    ],
    top=Inches(1.4),
    col_widths=[3.5, 4.5, 4.5],
)
add_body(sl, [
    "🔑 Puertos conocidos: HTTP=80  |  HTTPS=443  |  FTP=21  |  DNS=53  |  SMTP=25",
], top=Inches(5.3), font_size=Pt(17))

# =============================================================================
# SLIDE 8 – Modelo TCP/IP: Las 4 Capas
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "🌍 Modelo TCP/IP – Las 4 Capas")
add_table(sl,
    ["Capa TCP/IP", "Equivalente OSI", "Protocolos principales"],
    [
        ["Aplicación",    "Capas 5, 6, 7", "HTTP, FTP, DNS, SMTP, Telnet"],
        ["Transporte",    "Capa 4",        "TCP, UDP"],
        ["Internet",      "Capa 3",        "IP, ICMP, ARP, OSPF"],
        ["Acceso a Red",  "Capas 1 y 2",   "Ethernet, WiFi (802.11), PPP"],
    ],
    top=Inches(1.5),
    col_widths=[3.2, 3.5, 5.8],
)
add_body(sl, [
    "📌 TCP/IP = modelo PRÁCTICO — es el que realmente usa Internet hoy en día.",
    "📌 OSI = modelo TEÓRICO de referencia para entender y diseñar redes.",
], top=Inches(5.1), font_size=Pt(17))

# =============================================================================
# SLIDE 9 – OSI vs TCP/IP: Comparación Visual
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "🔄 OSI vs TCP/IP – Comparación")
add_table(sl,
    ["Capa OSI", "N°", "Capa TCP/IP"],
    [
        ["Aplicación",      "7", "Aplicación"],
        ["Presentación",    "6", "Aplicación"],
        ["Sesión",          "5", "Aplicación"],
        ["Transporte",      "4", "Transporte"],
        ["Red",             "3", "Internet"],
        ["Enlace de Datos", "2", "Acceso a Red"],
        ["Física",          "1", "Acceso a Red"],
    ],
    top=Inches(1.4),
    col_widths=[4.5, 1.0, 7.0],
)
add_body(sl, [
    '💡 "OSI para ENTENDER — TCP/IP para IMPLEMENTAR. Ambos modelos coexisten."',
], top=Inches(5.3), font_size=Pt(17))

# =============================================================================
# SLIDE 10 – Ejemplo Práctico: Navegar www.google.com
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "💻 Ejemplo Práctico: Abrir www.google.com")
add_body(sl, [
    "🔽 Proceso al hacer clic en el navegador (ocurre en milisegundos):",
    "",
    "  7️⃣ Aplicación  → Chrome genera petición HTTP GET /",
    "  4️⃣ Transporte  → TCP divide datos en segmentos, usa puerto 80",
    "  3️⃣ Red         → IP agrega dirección destino 142.250.x.x (Google)",
    "  2️⃣ Enlace      → Trama Ethernet con MAC address del router",
    "  1️⃣ Física      → Bits viajan por el cable UTP o WiFi",
    "",
    "  📨 Google responde y el proceso ocurre al revés (subiendo capas).",
    "",
    '  ✅ "Cada vez que abres una página, las 7 capas trabajan en conjunto."',
])

# =============================================================================
# SLIDE 11 – El Router
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "🔌 El Router – Conecta Redes Distintas")
add_body(sl, [
    "❓ ¿Qué es?: Dispositivo que CONECTA redes distintas y decide la mejor ruta para los paquetes.",
    "",
    "🔵 Capa OSI: Capa 3 – Red",
    "",
    "🏠 Tipos de router:",
    "    • Doméstico: router de casa con WiFi (ej. Movistar router)",
    "    • Empresarial: Cisco, Mikrotik — múltiples interfaces y VLANs",
    "    • De borde (ISP): conecta la red del proveedor con Internet",
    "",
    '💡 Ejemplo: "El router de tu casa tiene DOS interfaces:',
    '   → WAN (hacia Internet)   y   → LAN (hacia tu red local)."',
])

# =============================================================================
# SLIDE 12 – El Switch
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "🔀 El Switch – Conecta Equipos en la LAN")
add_body(sl, [
    "❓ ¿Qué es?: Dispositivo que conecta equipos DENTRO de la misma red (LAN).",
    "",
    "🔵 Capa OSI: Capa 2 – Enlace de Datos",
    "",
    "📋 Función: Aprende la MAC de cada equipo y reenvía tramas SOLO al destinatario.",
    "",
    "🔧 Tipos:",
    "    • No administrable: plug & play, sin configuración",
    "    • Administrable: soporta VLANs, QoS, trunking (Cisco Catalyst)",
    "",
    "⚖️  Switch vs Hub:",
    "    • Switch → envía trama SOLO al destino correcto (eficiente)",
    "    • Hub    → envía trama a TODOS los puertos (obsoleto)",
    "",
    '💡 Ejemplo: "El switch del laboratorio conecta los 30 PC entre sí."',
])

# =============================================================================
# SLIDE 13 – Patch Panel
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "🧩 Patch Panel – Centro Organizador del Rack")
add_body(sl, [
    "❓ ¿Qué es?: Panel con puertos RJ-45 en el rack que centraliza TODOS los cables del edificio.",
    "",
    "🎯 Función: Punto de conexión organizado entre el cableado permanente (pared) y los equipos activos.",
    "",
    "✅ Ventaja clave: Cambiar conexiones usando patch cords SIN mover cables de la pared.",
    "",
    "📦 Capacidades típicas: 24 puertos o 48 puertos (en 1U de rack).",
    "",
    "🔧 Se termina con herramienta punch-down (IDC) en la parte trasera.",
    "",
    '💡 Ejemplo: "El técnico desconecta el patch cord del patch panel para reubicar',
    '   un equipo — no toca los cables que están dentro de la pared."',
])

# =============================================================================
# SLIDE 14 – Jack / Wall Mount (Roseta de Red)
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "🔲 Jack / Roseta de Red – Acceso del Usuario")
add_body(sl, [
    "❓ ¿Qué es?: Conector hembra RJ-45 instalado en la pared.",
    "",
    "🏷️  También llamado: roseta, toma de red, faceplate, wall outlet.",
    "",
    "🎯 Función: Punto de acceso del usuario para conectar su equipo a la red.",
    "",
    "🔧 Instalación: Se termina con herramienta de impacto (punch-down tool).",
    "",
    "📏 Estándar de terminación: TIA/EIA 568A  o  TIA/EIA 568B.",
    "",
    "🎨 Código de colores: normalmente azul (datos) o negro/gris (voz).",
    "",
    '💡 Ejemplo: "El conector azul en la pared de la sala de clases donde enchufas',
    '   tu cable de red es un jack RJ-45 terminado en 568B."',
])

# =============================================================================
# SLIDE 15 – El Rack
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "🗄️ El Rack – Gabinete de Equipos de Red")
add_body(sl, [
    "❓ ¿Qué es?: Gabinete metálico estandarizado para montar equipos de red en forma ordenada.",
    "",
    "📐 Medida estándar: Unidades de Rack (U) → 1U = 44.45 mm de altura.",
    "",
    "📦 Tamaños comunes: 12U (pequeño) | 24U (mediano) | 42U (sala de servidores).",
    "",
    "🔧 Contenido típico de un rack:",
    "    • Switches administrables",
    "    • Patch panels (24 o 48 puertos)",
    "    • Servidores 1U / 2U",
    "    • UPS (UPS de rack)",
    "    • Organizadores de cable (1U)",
    "",
    '💡 Ejemplo: "Un rack de 42U puede contener hasta 42 dispositivos de 1U apilados."',
])

# =============================================================================
# SLIDE 16 – Servidores
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "🖥️ Servidores – Proveedores de Servicios en Red")
add_body(sl, [
    "❓ ¿Qué es?: Computador de alta disponibilidad que provee servicios a clientes de red.",
    "",
    "⚖️  Diferencia con un PC normal: hardware redundante, sin pantalla, diseño 1U/2U para rack.",
    "",
    "🔧 Tipos de servidores y su función:",
    "    • DHCP    → Asigna IPs automáticamente a los equipos de la red",
    "    • DNS     → Traduce nombres de dominio a direcciones IP",
    "    • Web     → Aloja páginas web (Apache, IIS, Nginx)",
    "    • Archivos → Almacenamiento compartido (NAS, Samba)",
    "    • AD/LDAP → Gestión de usuarios y permisos (Active Directory)",
    "",
    '💡 Ejemplo: "El servidor DHCP del colegio le da una IP a tu notebook cuando',
    '   te conectas al WiFi — automáticamente, en menos de 1 segundo."',
])

# =============================================================================
# SLIDE 17 – Diagrama: Cómo se conectan los componentes
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "🗺️ Diagrama: Flujo de Conexión de la Red")
add_body(sl, [
    "  🌐 INTERNET",
    "       ↓",
    "  📡 ACOMETIDA (fibra del ISP) ──→ 🖥️ ROUTER (Capa 3)",
    "                                          ↓",
    "                                   🔀 SWITCH CORE (Capa 2)",
    "                                    ↙             ↘",
    "              🔀 Switch Acceso Piso 1         🖥️ Servidores (DHCP, DNS, Web)",
    "                      ↓",
    "              🧩 PATCH PANEL (rack)",
    "                      ↓",
    "              〰️  Cable UTP Cat6 por la pared",
    "                      ↓",
    "              🔲 JACK RJ-45 (en la pared)",
    "                      ↓",
    "              💻 PC del usuario (con patch cord)",
])

# =============================================================================
# SLIDE 18 – Resumen Visual: Dispositivos y su Capa OSI
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "📊 Resumen: Dispositivos y su Capa OSI")
add_table(sl,
    ["Dispositivo", "Capa OSI", "Función principal"],
    [
        ["Hub",          "Capa 1 – Física",          "Repite señal a TODOS los puertos (obsoleto)"],
        ["Switch",       "Capa 2 – Enlace de Datos",  "Reenvía tramas por dirección MAC"],
        ["Router",       "Capa 3 – Red",              "Enruta paquetes por dirección IP"],
        ["Patch Panel",  "N/A (pasivo)",               "Organiza y centraliza el cableado"],
        ["Jack/Roseta",  "N/A (pasivo)",               "Punto de acceso del usuario a la red"],
        ["Servidor",     "Capa 7 – Aplicación",        "Provee servicios: DHCP, DNS, Web, etc."],
    ],
    top=Inches(1.4),
    col_widths=[2.5, 3.0, 7.0],
)

# =============================================================================
# SLIDE 19 – ¿Qué es el Cableado Estructurado?
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "🔧 ¿Qué es el Cableado Estructurado?")
add_body(sl, [
    "📌 Definición: Sistema de cableado organizado, estandarizado y documentado para transmitir",
    "   datos, voz y video en un edificio o campus.",
    "",
    "📏 Norma que lo regula: TIA/EIA-568 (internacional).",
    "",
    "🔑 Características clave:",
    "    • Independiente del fabricante → cualquier marca de equipo es compatible",
    "    • Escalable → se puede crecer la red sin rehacer el cableado",
    "    • Documentado → planos, etiquetas e inventario de cada punto",
    "",
    "✅ Por qué es importante:",
    "    ✅ Reduce fallas y tiempo de reparación",
    "    ✅ Permite crecer la red fácilmente",
    "    ✅ Acepta datos, voz y video en la misma infraestructura",
    "    ✅ Exigido en instalaciones profesionales y licitaciones",
])

# =============================================================================
# SLIDE 20 – Importancia del Cableado Estructurado
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "⚠️ ¿Por qué es tan Importante?")
add_body(sl, [
    "❌ SIN cableado estructurado:",
    "    • Cables sin identificar → nadie sabe a dónde van",
    "    • Difícil mantención → reparaciones toman horas o días",
    "    • Aspecto caótico → aumenta riesgo de errores",
    "",
    "✅ CON cableado estructurado:",
    "    • Cada cable etiquetado y documentado",
    "    • Rutas ordenadas por canaletas y bandejas",
    "    • Plano actualizado de cada punto de red",
    "",
    "📊 Dato real: El 70% de los problemas de red tienen origen en el cableado mal instalado.",
    "",
    '💡 "Un buen cableado estructurado puede durar 15–20 años aunque cambies',
    '   todos los equipos activos de la red."',
])

# =============================================================================
# SLIDE 21 – El Backbone (Cableado Vertical)
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "🏗️ El Backbone – Cableado Vertical")
add_body(sl, [
    "❓ ¿Qué es?: Cableado que interconecta los diferentes pisos, edificios o salas de telecomunicaciones.",
    "",
    "🏷️  También llamado: cableado vertical, riser, columna vertebral.",
    "",
    "📡 Medio usado: Fibra Óptica (preferido) o UTP Cat6A para distancias cortas.",
    "",
    "🎯 Función: Transportar el tráfico de TODA la red hacia el punto central (MDF).",
    "",
    "🏢 En un edificio de 3 pisos:",
    "    Sala principal (MDF) → Fibra Óptica → Sala de piso IDF → UTP → Usuario",
    "",
    '🛣️  Analogía: "El backbone es como la AUTOPISTA de la red.',
    '   Los cables de piso son las calles locales del barrio."',
])

# =============================================================================
# SLIDE 22 – Sala de Servidores / Cuarto de Telecomunicaciones
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "🏢 MDF e IDF – Salas de Telecomunicaciones")
add_body(sl, [
    "🖥️  MDF – Main Distribution Frame (Sala de Servidores Principal):",
    "    • Centro PRINCIPAL de toda la red del edificio",
    "    • Contiene: router, switch core, servidores, UPS principal",
    "    • Requisitos: climatización (A/C), acceso restringido, piso técnico",
    "",
    "📦 IDF – Intermediate Distribution Frame (Cuarto de Telecomunicaciones):",
    "    • Sala SECUNDARIA por piso o área del edificio",
    "    • Contiene: switch de acceso, patch panel de ese piso",
    "    • Conectado al MDF mediante backbone de fibra óptica",
    "",
    '💡 Ejemplo en un colegio de 3 pisos:',
    '   → 1 MDF en planta baja (sala de computación principal)',
    '   → 1 IDF por piso adicional (armario de red en pasillo)',
])

# =============================================================================
# SLIDE 23 – La Acometida
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "📡 La Acometida – La Conexión que viene de Afuera")
add_body(sl, [
    "❓ ¿Qué es la Acometida?: Tramo de red que va desde la infraestructura pública",
    "   (poste o cámara de calle) hasta el interior del edificio.",
    "",
    "🏷️  También conocida como: 'última milla', 'acometida externa', 'ingreso de portador'.",
    "",
    "👷 ¿Quién la instala?: El proveedor de Internet (ISP) — Movistar, VTR, Entel, GTD, etc.",
    "",
    "📡 Medios usados según tecnología:",
    "    • Fibra Óptica (FTTH)    → mejor velocidad y calidad (hasta 10 Gbps)",
    "    • Cable coaxial (HFC)    → TV cable + Internet",
    "    • Par de cobre (ADSL)    → tecnología antigua, menor velocidad",
    "",
    "🔌 Termina en: Sala de Servidores (MDF) en el equipo ONT o módem del ISP.",
    "",
    '💡 Ejemplo: "El cable naranja que baja del poste y entra por el muro del colegio',
    '   es la acometida de Movistar — desde ahí en adelante es red interna."',
])

# =============================================================================
# SLIDE 24 – Importancia de Marcar los Cables
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "🏷️ Por qué es Importante Marcar los Cables")
add_body(sl, [
    "📏 Norma de referencia: TIA-606 — regula la administración e identificación del cableado.",
    "",
    "🔑 Razones para marcar los cables:",
    "    🔍 Identificar un cable en segundos, no en horas",
    "    📋 Mantener documentación actualizada de la red",
    "    🔧 Reducir errores en mantención (no desconectar el cable equivocado)",
    "    👷 Requerido para certificaciones, auditorías y garantías",
    "",
    "🏷️  Métodos de marcado:",
    "    • Etiquetas adhesivas impresas en AMBOS extremos del cable",
    "    • Código de colores por área, piso o servicio",
    "    • Numeración estándar: PISO-SALA-PUNTO  (ej: 2-LAB-PC05)",
    "",
    '⚡ Ejemplo real:',
    '   Sin marcado → técnico tarda 2 HORAS buscando el cable del PC-15',
    '   Con marcado → encuentra el cable en 30 SEGUNDOS',
])

# =============================================================================
# SLIDE 25 – Subsistemas del Cableado Estructurado
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "📐 Subsistemas del Cableado Estructurado (TIA-568)")
add_table(sl,
    ["Subsistema", "Descripción"],
    [
        ["1. Entrada del Edificio",    "Acometida externa → sala de servidores (MDF)"],
        ["2. Sala de Equipos",         "MDF con servidores, router y switch core"],
        ["3. Cableado Backbone",       "Conecta MDF con IDFs → usa fibra óptica"],
        ["4. Cuarto de Telecom",       "IDF por piso: patch panel + switch de acceso"],
        ["5. Cableado Horizontal",     "Del patch panel al jack de pared → UTP Cat6 (máx 90m)"],
        ["6. Área de Trabajo",         "Del jack RJ-45 al equipo del usuario → patch cord"],
    ],
    top=Inches(1.4),
    col_widths=[3.8, 8.7],
)

# =============================================================================
# SLIDE 26 – ¿Qué son las Normas TIA/EIA-568?
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "📏 Normas TIA/EIA-568 – ¿Para qué sirven?")
add_body(sl, [
    "🏢 Publicadas por: TIA (Telecommunications Industry Association) y",
    "   EIA (Electronic Industries Alliance) — USA.",
    "",
    "🎯 ¿Para qué sirven?: Estandarizar el cableado de telecomunicaciones en edificios.",
    "",
    "📋 Incluyen: tipos de cable aceptados, distancias máximas, conectores y ORDEN DE COLORES.",
    "",
    "🗓️  Versiones históricas:",
    "    • TIA-568-A → 1991 (primera versión)",
    "    • TIA-568-B → 2001 (más usada en Chile)",
    "    • TIA-568-C → 2009 (actualización)",
    "    • TIA-568-D → 2015 (incluye Cat8)",
    "",
    '⚠️  "Sin estas normas, cada empresa haría el cableado a su manera',
    '   y nada sería compatible entre sí."',
])

# =============================================================================
# SLIDE 27 – El Conector RJ-45
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "🔌 El Conector RJ-45 – 8 Pines, 8 Conductores")
add_body(sl, [
    "📌 RJ-45: 8P8C → 8 Posiciones / 8 Conductores (contactos).",
    "",
    "🔧 Partes del conector:",
    "    • Cuerpo plástico transparente",
    "    • 8 pines dorados de contacto",
    "    • Lengüeta de seguridad (clip) para bloquear el conector",
    "",
    "🛠️  Herramienta para instalar: PONCHADORA RJ-45",
    "",
    "📡 Cables compatibles: UTP Cat5e, Cat6, Cat6A",
    "",
    "🎯 Usado en: Ethernet (redes de datos LAN)",
    "",
    '⚠️  TIP importante: "Siempre verificar el ORDEN DE COLORES ANTES de ponchar.',
    '   Un error en el orden significa rehacer el conector completo."',
])

# =============================================================================
# SLIDE 28 – Orden de Colores TIA 568A
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "🟢 Orden de Colores – TIA 568A")
add_table(sl,
    ["Pin", "Color TIA 568A", "Par"],
    [
        ["1", "⚪ Blanco / Verde",   "Par 3"],
        ["2", "🟢 Verde",            "Par 3"],
        ["3", "⚪ Blanco / Naranja", "Par 2"],
        ["4", "🔵 Azul",             "Par 1"],
        ["5", "⚪ Blanco / Azul",    "Par 1"],
        ["6", "🟠 Naranja",          "Par 2"],
        ["7", "⚪ Blanco / Café",    "Par 4"],
        ["8", "🟤 Café",             "Par 4"],
    ],
    top=Inches(1.4),
    col_widths=[1.0, 7.0, 4.5],
)
add_body(sl, [
    "🧠 Nemotecnia 568A: BV – V – BN – Az – BAz – N – BCa – Ca",
], top=Inches(5.5), font_size=Pt(17))

# =============================================================================
# SLIDE 29 – Orden de Colores TIA 568B
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "🟠 Orden de Colores – TIA 568B")
add_table(sl,
    ["Pin", "Color TIA 568B", "Par"],
    [
        ["1", "⚪ Blanco / Naranja", "Par 2"],
        ["2", "🟠 Naranja",          "Par 2"],
        ["3", "⚪ Blanco / Verde",   "Par 3"],
        ["4", "🔵 Azul",             "Par 1"],
        ["5", "⚪ Blanco / Azul",    "Par 1"],
        ["6", "🟢 Verde",            "Par 3"],
        ["7", "⚪ Blanco / Café",    "Par 4"],
        ["8", "🟤 Café",             "Par 4"],
    ],
    top=Inches(1.4),
    col_widths=[1.0, 7.0, 4.5],
)
add_body(sl, [
    "🧠 Nemotecnia 568B: BN – N – BV – Az – BAz – V – BCa – Ca",
    "⚠️  Diferencia con 568A: solo intercambia los pares Naranja (pins 1,2) y Verde (pins 3,6).",
], top=Inches(5.3), font_size=Pt(16))

# =============================================================================
# SLIDE 30 – Cable Directo vs Cable Cruzado
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "🔁 Cable Directo vs Cable Cruzado")
add_table(sl,
    ["Característica", "Cable Directo", "Cable Cruzado"],
    [
        ["Estándar",          "568B – 568B (mismo en ambos extremos)", "568A – 568B (distintos extremos)"],
        ["Color pin 1",       "⚪ Blanco / Naranja",                    "⚪ Blanco / Verde"],
        ["Equipos conectados","Equipos DISTINTOS",                     "Equipos IGUALES"],
        ["Uso típico",        "PC → Switch | Switch → Router",         "PC → PC | Switch → Switch"],
        ["¿Se usa hoy?",      "✅ Sí, instalaciones estándar",          "⚠️ Raro — Auto MDI-X lo reemplaza"],
    ],
    top=Inches(1.4),
    col_widths=[3.0, 5.0, 5.0],
)
add_body(sl, [
    "🤖 Auto MDI-X: tecnología moderna que detecta el tipo de conexión automáticamente.",
    '✅ "En instalaciones nuevas siempre se usa 568B en ambos extremos (cable directo)."',
], top=Inches(5.3), font_size=Pt(16))

# =============================================================================
# SLIDE 31 – Resumen Final y Buenas Prácticas
# =============================================================================
sl = add_slide()
set_background(sl)
orange_bar(sl)
add_title(sl, "✅ Resumen Final y Buenas Prácticas")
add_body(sl, [
    "✅ Usar SIEMPRE el mismo estándar en toda la instalación (generalmente 568B).",
    "✅ Marcar AMBOS extremos del cable con la misma etiqueta.",
    "✅ Respetar distancias máximas: 90m cable permanente + 10m patch cords = 100m total.",
    "✅ Documentar cada punto de red en un plano actualizado.",
    "✅ Usar materiales certificados (cables con etiqueta TIA/EIA).",
    "",
    "⚠️  NO mezclar 568A y 568B en la misma instalación.",
    "⚠️  NO doblar cables en ángulos menores a 4× el diámetro del cable.",
    "⚠️  NO superar los 100m de distancia total por punto de red.",
    "",
    '🏆 "Una instalación bien hecha HOY evita horas de trabajo MAÑANA.',
    '   El cableado estructurado es la base de toda red profesional."',
])

# =============================================================================
# Guardar presentación
# =============================================================================
output = "resumen_modulos.pptx"
prs.save(output)
print(f"✅ Presentación guardada: {output}")
print(f"   Total de diapositivas: {len(prs.slides)}")
